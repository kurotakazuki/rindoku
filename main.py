import argparse
from openai import OpenAI
from tqdm import tqdm
from datetime import datetime
import os
import json
from pptx import Presentation
from pptx.util import Cm


class JsonToPptx:
    def __init__(self):
        self.ppt = Presentation()
        JsonToPptx.update_slide_layout(self.ppt)

    def update_slide_layout(ppt):
        ppt.slide_width = Cm(33.867)
        ppt.slide_height = Cm(19.05)

    def add_content(slide, content, level=0):
        for item in content:
            text = item["text"]
            p = slide.shapes.placeholders[1].text_frame.add_paragraph()
            p.text = text
            p.level = level
            if "childs" in item and item["childs"]:
                JsonToPptx.add_content(slide, item["childs"], level + 1)

    def create_pptx(self, json_txt, file_name="presentation.pptx"):
        json_data = json.loads(json_txt)
        for slide_data in json_data["slides"]:
            slide_layout = self.ppt.slide_layouts[1]
            slide = self.ppt.slides.add_slide(slide_layout)

            # Set slide title
            title = slide.shapes.title
            title.text = slide_data["title"]

            # Add slide content
            if "content" in slide_data:
                JsonToPptx.add_content(slide, slide_data["content"])

            # Add slide note
            if "original" in slide_data:
                slide.notes_slide.notes_text_frame.text = slide_data["original"]

        self.ppt.save(file_name)


class Rindoku:

    def __init__(self):
        print("Rindoku CLI INIT")
        args = Rindoku.get_args()
        self.prompt_str = Rindoku.read_prompt(args.prompt)
        self.model = args.model
        self.input = args.input
        self.json_prompt_str = Rindoku.read_json_prompt(args.json_prompt)
        self.json2pptx = JsonToPptx()

    def get_args():
        print("Rindoku CLI GET_ARGS")
        parser = argparse.ArgumentParser(description="Rindoku CLI")
        parser.add_argument(
            "-i",
            "--input",
            type=str,
            default="",
            help="Input text",
        )
        parser.add_argument(
            "-m", "--model", type=str, default="gpt-4o", help="Model name"
        )
        parser.add_argument(
            "-p", "--prompt", type=str, default="prompt.txt", help="Prompt file path"
        )
        parser.add_argument(
            "-j",
            "--json-prompt",
            type=str,
            default="json_prompt.txt",
            help="JSON Prompt file path",
        )
        return parser.parse_args()

    def read_prompt(prompt_path):
        print("Rindoku CLI READ_PLOT_PROMPT")
        try:
            with open(prompt_path, "r", encoding="utf-8") as file:
                file_content = file.read()
        except Exception as e:
            print(
                f"Error reading the plot prompt file: {e}\nCreate an original prompt and set it in the `prompt.txt` file."
            )
            exit(1)
        return file_content

    def read_json_prompt(prompt_path):
        print("Rindoku CLI READ_JSON_PROMPT")
        try:
            with open(prompt_path, "r", encoding="utf-8") as file:
                file_content = file.read()
        except Exception as e:
            print(
                f"Error reading the json prompt file: {e}\nCreate an original prompt and set it in the `json_prompt.txt` file."
            )
            exit(1)
        return file_content

    def file_name(folder_name, file_type="", extension="txt"):
        current_time = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
        return f"{folder_name}/{current_time}_{file_type}.{extension}"

    def save_string_to_file(output_string, file_type="", extension="txt"):
        folder_name = "output"
        if not os.path.exists(folder_name):
            os.makedirs(folder_name)

        file_name = Rindoku.file_name(folder_name, file_type, extension)

        with open(file_name, "w", encoding="utf-8") as file:
            file.write(output_string)

    def create_plot(self, client, input):
        completion = client.chat.completions.create(
            model=self.model,
            messages=[
                {
                    "role": "system",
                    "content": self.prompt_str,
                },
                {"role": "user", "content": input},
            ],
        )

        slide_plot = completion.choices[0].message.content
        print(slide_plot)
        Rindoku.save_string_to_file(slide_plot, file_type="plot")
        return slide_plot

    def create_json(self, client, slide_plot):
        completion = client.chat.completions.create(
            model=self.model,
            messages=[
                {
                    "role": "system",
                    "content": self.json_prompt_str,
                },
                {"role": "user", "content": slide_plot},
            ],
        )

        json_str = completion.choices[0].message.content
        json_str = json_str.replace("```json", "")
        json_str = json_str.replace("```", "")
        print(json_str)
        Rindoku.save_string_to_file(json_str, file_type="json", extension="json")
        return json_str

    def run(self):
        print("Rindoku CLI RUN")
        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            print(
                "Error: OPENAI_API_KEY is not set in the environment variables.\nCreate an API key and set it in the `.env` file that references `.env.example`."
            )
            exit(1)

        client = OpenAI(api_key=api_key)

        print("Rindoku CLI create slide plot")
        slide_plot = self.create_plot(client, self.input)

        print("Rindoku CLI create slide json")
        json_str = self.create_json(client, slide_plot)

        self.json2pptx.create_pptx(
            json_str,
            file_name=Rindoku.file_name(
                "output", file_type="presentation", extension="pptx"
            ),
        )


rindoku = Rindoku()
rindoku.run()
